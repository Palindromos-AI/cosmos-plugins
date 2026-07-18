from __future__ import annotations

import contextlib
import importlib.util
import io
import json
from pathlib import Path
import sys
import tempfile
from types import SimpleNamespace
import unittest
from unittest import mock
import zipfile


SCRIPT = Path(__file__).parents[1] / "scripts" / "raw_to_markdown.py"
SPEC = importlib.util.spec_from_file_location("raw_to_markdown", SCRIPT)
assert SPEC and SPEC.loader
MODULE = importlib.util.module_from_spec(SPEC)
sys.modules[SPEC.name] = MODULE
SPEC.loader.exec_module(MODULE)


class RawToMarkdownTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temp = tempfile.TemporaryDirectory()
        self.vault = Path(self.temp.name)
        self.raw = self.vault / "raw" / "dollar$" / "case"
        self.raw.mkdir(parents=True)

    def tearDown(self) -> None:
        self.temp.cleanup()

    def test_plan_creates_same_directory_sidecar(self) -> None:
        source = self.raw / "notes.txt"
        source.write_text("Alpha\nBeta\n", encoding="utf-8")

        items = MODULE.build_plan([str(source)], self.vault, recursive=False)

        self.assertEqual(items[0].action, "create")
        self.assertEqual(items[0].source, "raw/dollar$/case/notes.txt")
        self.assertEqual(items[0].output, "raw/dollar$/case/notes.md")

    def test_path_outside_raw_is_rejected(self) -> None:
        source = self.vault / "outside.txt"
        source.write_text("No", encoding="utf-8")

        items = MODULE.build_plan([str(source)], self.vault, recursive=False)

        self.assertEqual(items[0].action, "invalid")

    def test_external_directory_is_rejected_before_recursive_expansion(self) -> None:
        outside = self.vault / "outside"
        outside.mkdir()
        (outside / "a.txt").write_text("A", encoding="utf-8")
        (outside / "b.txt").write_text("B", encoding="utf-8")

        items = MODULE.build_plan([str(outside)], self.vault, recursive=True)

        self.assertEqual(len(items), 1)
        self.assertEqual(items[0].action, "invalid")
        self.assertEqual(items[0].source, str(outside))

    def test_plan_cli_reports_unsafe_selection_as_failure(self) -> None:
        outside = self.vault / "outside.txt"
        outside.write_text("No", encoding="utf-8")
        stdout = io.StringIO()

        with contextlib.redirect_stdout(stdout):
            exit_code = MODULE.main(
                ["plan", "--vault", str(self.vault), str(outside)]
            )

        report = json.loads(stdout.getvalue())
        self.assertEqual(exit_code, 2)
        self.assertFalse(report["ok"])
        self.assertEqual(report["items"][0]["action"], "invalid")

    def test_existing_unowned_markdown_is_never_overwritten(self) -> None:
        source = self.raw / "report.txt"
        output = self.raw / "report.md"
        source.write_text("Source", encoding="utf-8")
        output.write_text("Human note", encoding="utf-8")

        items = MODULE.build_plan([str(source)], self.vault, recursive=False)

        self.assertEqual(items[0].action, "conflict")
        self.assertEqual(output.read_text(encoding="utf-8"), "Human note")

    def test_markdown_created_during_conversion_is_never_overwritten(self) -> None:
        source = self.raw / "race.txt"
        output = source.with_suffix(".md")
        source.write_text("Source", encoding="utf-8")

        class RacingEngine:
            def convert_local(self, _source: Path) -> SimpleNamespace:
                output.write_text("Human note created concurrently", encoding="utf-8")
                return SimpleNamespace(text_content="Generated body")

        with (
            mock.patch.object(MODULE, "load_engine", return_value=RacingEngine()),
            mock.patch.object(
                MODULE, "engine_version", return_value=MODULE.PINNED_ENGINE_VERSION
            ),
        ):
            results, exit_code = MODULE.run_convert(
                [str(source)], self.vault, recursive=False
            )

        self.assertEqual(exit_code, 1)
        self.assertEqual(results[0]["action"], "failed")
        self.assertEqual(
            output.read_text(encoding="utf-8"), "Human note created concurrently"
        )
        self.assertFalse((output.parent / f".{output.name}.raw-to-markdown.lock").exists())

    def test_convert_preserves_source_and_writes_provenance(self) -> None:
        source = self.raw / "source.txt"
        original = b"Heading\n\nBody text.\n"
        source.write_bytes(original)

        results, exit_code = MODULE.run_convert(
            [str(source)], self.vault, recursive=False
        )

        self.assertEqual(exit_code, 0)
        self.assertEqual(results[0]["action"], "created")
        self.assertEqual(source.read_bytes(), original)
        output = source.with_suffix(".md")
        content = output.read_text(encoding="utf-8")
        metadata, body = MODULE.split_frontmatter(content)
        self.assertEqual(metadata["conversion_schema"], "raw-to-markdown/v2")
        self.assertEqual(metadata["converted_from_path"], "raw/dollar$/case/source.txt")
        self.assertEqual(metadata["converted_from_sha256"], MODULE.sha256_bytes(original))
        self.assertEqual(metadata["conversion_postprocessor"], "none")
        self.assertIn("[[raw/dollar$/case/source.txt]]", body)
        self.assertIn("Body text.", body)
        self.assertEqual(
            MODULE.classify_source(source, "raw/dollar$/case/source.txt").action,
            "no-op",
        )

    def test_changed_source_marks_owned_sidecar_stale_without_overwriting(self) -> None:
        source = self.raw / "update.txt"
        source.write_text("Version one", encoding="utf-8")
        _, first_exit = MODULE.run_convert(
            [str(source)], self.vault, recursive=False
        )
        self.assertEqual(first_exit, 0)
        output = source.with_suffix(".md")
        previous = output.read_bytes()
        source.write_text("Version two", encoding="utf-8")
        self.assertEqual(
            MODULE.classify_source(source, "raw/dollar$/case/update.txt").action,
            "stale-conflict",
        )

        results, second_exit = MODULE.run_convert(
            [str(source)], self.vault, recursive=False
        )

        self.assertEqual(second_exit, 2)
        self.assertEqual(results[0]["action"], "stale-conflict")
        self.assertEqual(output.read_bytes(), previous)

    def test_manual_sidecar_edit_is_never_overwritten(self) -> None:
        source = self.raw / "edited.txt"
        source.write_text("Source", encoding="utf-8")
        _, exit_code = MODULE.run_convert(
            [str(source)], self.vault, recursive=False
        )
        self.assertEqual(exit_code, 0)
        output = source.with_suffix(".md")
        output.write_text(output.read_text(encoding="utf-8") + "Manual edit\n", encoding="utf-8")

        blocked, blocked_exit = MODULE.run_convert(
            [str(source)], self.vault, recursive=False
        )

        self.assertEqual(blocked_exit, 2)
        self.assertEqual(blocked[0]["action"], "edited-conflict")
        self.assertIn("Manual edit", output.read_text(encoding="utf-8"))

    def test_missing_provenance_field_is_an_unrelated_file_conflict(self) -> None:
        source = self.raw / "provenance.txt"
        source.write_text("Source", encoding="utf-8")
        _, exit_code = MODULE.run_convert(
            [str(source)], self.vault, recursive=False
        )
        self.assertEqual(exit_code, 0)
        output = source.with_suffix(".md")
        lines = output.read_text(encoding="utf-8").splitlines(keepends=True)
        output.write_text(
            "".join(line for line in lines if not line.startswith("conversion_engine: ")),
            encoding="utf-8",
        )

        items = MODULE.build_plan([str(source)], self.vault, recursive=False)
        blocked, blocked_exit = MODULE.run_convert(
            [str(source)], self.vault, recursive=False
        )

        self.assertEqual(items[0].action, "conflict")
        self.assertEqual(blocked_exit, 2)
        self.assertEqual(blocked[0]["action"], "conflict")

    def test_converter_failure_leaves_no_sidecar_or_lock(self) -> None:
        source = self.raw / "broken.txt"
        output = source.with_suffix(".md")
        source.write_text("Source", encoding="utf-8")

        class BrokenEngine:
            def convert_local(self, _source: Path) -> SimpleNamespace:
                raise RuntimeError("converter failed")

        with (
            mock.patch.object(MODULE, "load_engine", return_value=BrokenEngine()),
            mock.patch.object(
                MODULE, "engine_version", return_value=MODULE.PINNED_ENGINE_VERSION
            ),
        ):
            results, exit_code = MODULE.run_convert(
                [str(source)], self.vault, recursive=False
            )

        self.assertEqual(exit_code, 1)
        self.assertEqual(results[0]["action"], "failed")
        self.assertFalse(output.exists())
        self.assertFalse((output.parent / f".{output.name}.raw-to-markdown.lock").exists())

    def test_pathological_pdf_table_detection_ignores_small_real_table(self) -> None:
        small_table = "\n".join(
            ["| Metric | Value |", "| --- | --- |", "| Carbon | 42 |"]
        )
        prose = "这是被错误拆入表格的连续中文正文" * 4
        pathological = "\n".join(
            line
            for _ in range(10)
            for line in (f"| {prose} |  |  |", "| --- | --- | --- |")
        )

        self.assertFalse(MODULE.has_pathological_pdf_tables(small_table))
        self.assertTrue(MODULE.has_pathological_pdf_tables(pathological))

    def test_pdf_prose_cleanup_removes_margins_page_numbers_and_wraps(self) -> None:
        cleaned = MODULE.clean_pdf_pages(
            [
                "\n".join(
                    [
                        "2021年 北 京 工 业 大 学 学 报",
                        "中国实现 2030 年前碳达峰目标及主要途径",
                        "一、主要路径",
                        "中国实现2030 年前碳达",
                        "峰目标。",
                        "8",
                    ]
                ),
                "\n".join(
                    [
                        "2022年 北 京 工 业 大 学 学 报",
                        "(一)控制能源消费总量",
                        "大幅度提高非化石能源",
                        "占一次能源消费比重。",
                        "9",
                    ]
                ),
            ],
            "中国实现2030年前碳达峰目标及主要途径",
        )

        self.assertIn("# 中国实现2030年前碳达峰目标及主要途径", cleaned)
        self.assertIn("## 一、主要路径", cleaned)
        self.assertIn("### (一)控制能源消费总量", cleaned)
        self.assertIn("中国实现2030年前碳达峰目标。", cleaned)
        self.assertIn("大幅度提高非化石能源占一次能源消费比重。", cleaned)
        self.assertNotIn("北京工业大学学报", cleaned)
        self.assertNotRegex(cleaned, r"(?m)^[89]$")

    def test_pdf_prose_postprocessor_is_recorded(self) -> None:
        source = self.raw / "journal.pdf"
        source.write_bytes(b"fake pdf bytes")
        prose = "这是被错误拆入表格的连续中文正文" * 4
        pathological = "\n".join(
            line
            for _ in range(10)
            for line in (f"| {prose} |  |  |", "| --- | --- | --- |")
        )

        class TableHeavyEngine:
            def convert_local(self, _source: Path) -> SimpleNamespace:
                return SimpleNamespace(text_content=pathological)

        with (
            mock.patch.object(MODULE, "load_engine", return_value=TableHeavyEngine()),
            mock.patch.object(
                MODULE, "engine_version", return_value=MODULE.PINNED_ENGINE_VERSION
            ),
            mock.patch.object(
                MODULE,
                "extract_pdf_prose",
                return_value="# Clean PDF\n\nClean PDF marker.",
            ) as extract,
        ):
            results, exit_code = MODULE.run_convert(
                [str(source)], self.vault, recursive=False
            )

        self.assertEqual(exit_code, 0, results)
        extract.assert_called_once()
        content = source.with_suffix(".md").read_text(encoding="utf-8")
        metadata, body = MODULE.split_frontmatter(content)
        self.assertEqual(metadata["conversion_postprocessor"], "pdf-prose-v1")
        self.assertIn("Clean PDF marker.", body)
        self.assertNotIn("| --- |", body)

    def test_pdf_cleanup_preserves_reference_issue_and_separates_metadata(self) -> None:
        cleaned = MODULE.clean_pdf_pages(
            [
                "\n".join(
                    [
                        "中图分类号: F120; F205 文献标志码: A",
                        "世界已经进入全球气候变化时代。",
                        "［12］作者．期刊名称，2019",
                        "(11): 31-41，54．",
                        "［16］作者．研究标题［J］．",
                        "期刊名称，2019，29(4): 1-9．",
                    ]
                )
            ],
            "unmatched-title",
        )

        self.assertIn(
            "中图分类号: F120; F205文献标志码: A\n\n世界已经进入全球气候变化时代。",
            cleaned,
        )
        self.assertIn("［12］作者．期刊名称，2019(11): 31-41，54．", cleaned)
        self.assertIn(
            "［16］作者．研究标题［J］．期刊名称，2019，29(4): 1-9．", cleaned
        )
        self.assertNotIn("### (11)", cleaned)

    def test_pdf_cleanup_does_not_treat_decimal_continuation_as_affiliation(self) -> None:
        cleaned = MODULE.clean_pdf_pages(
            ["相当于美国比重\n(7.9%) 的6.0倍。\n排放量为\n(1.10亿吨)。"],
            "unmatched-title",
        )

        self.assertIn("相当于美国比重(7.9%) 的6.0倍。", cleaned)
        self.assertIn("排放量为(1.10亿吨)。", cleaned)

    def test_pdf_cleanup_moves_page_footnotes_and_joins_cross_page_prose(self) -> None:
        cleaned = MODULE.clean_pdf_pages(
            [
                "\n".join(
                    [
                        "2021年 北京工业大学学报",
                        "正文第一句。",
                        "正文第二句。",
                        "可再生能源总",
                        "① 脚注内容 https:∥example．org。",
                        "8",
                    ]
                ),
                "\n".join(
                    [
                        "2022年 北京工业大学学报",
                        "消费量占比提高。",
                        "9",
                    ]
                ),
            ],
            "unmatched-title",
        )

        self.assertIn("可再生能源总消费量占比提高。", cleaned)
        self.assertIn("## 注释\n\n① 脚注内容 https:∥example．org。", cleaned)
        self.assertLess(cleaned.index("消费量占比提高。"), cleaned.index("## 注释"))

    def test_source_change_after_create_preserves_sidecar_and_reports_failure(self) -> None:
        source = self.raw / "post-write-source-change.txt"
        output = source.with_suffix(".md")
        source.write_text("Version one", encoding="utf-8")
        real_create = MODULE.create_output

        def create_then_change_source(path: Path, content: bytes) -> None:
            real_create(path, content)
            source.write_text("Changed after commit", encoding="utf-8")

        with mock.patch.object(
            MODULE, "create_output", side_effect=create_then_change_source
        ):
            results, exit_code = MODULE.run_convert(
                [str(source)], self.vault, recursive=False
            )

        self.assertEqual(exit_code, 1)
        self.assertEqual(results[0]["action"], "failed")
        self.assertIn("output was preserved", results[0]["reason"])
        self.assertTrue(output.exists())
        self.assertEqual(
            MODULE.classify_source(source, "raw/dollar$/case/post-write-source-change.txt").action,
            "stale-conflict",
        )
        self.assertFalse((output.parent / f".{output.name}.raw-to-markdown.lock").exists())

    def test_concurrent_edit_after_create_is_preserved(self) -> None:
        source = self.raw / "post-write-edit.txt"
        output = source.with_suffix(".md")
        source.write_text("Source", encoding="utf-8")
        real_create = MODULE.create_output

        def create_then_edit(path: Path, content: bytes) -> None:
            real_create(path, content)
            output.write_text("HUMAN CONCURRENT EDIT", encoding="utf-8")

        with mock.patch.object(
            MODULE, "create_output", side_effect=create_then_edit
        ):
            results, exit_code = MODULE.run_convert(
                [str(source)], self.vault, recursive=False
            )

        self.assertEqual(exit_code, 1)
        self.assertEqual(results[0]["action"], "failed")
        self.assertIn("output was preserved", results[0]["reason"])
        self.assertEqual(output.read_text(encoding="utf-8"), "HUMAN CONCURRENT EDIT")
        self.assertFalse((output.parent / f".{output.name}.raw-to-markdown.lock").exists())

    def test_recursive_flag_is_required_for_directories(self) -> None:
        (self.raw / "a.txt").write_text("A", encoding="utf-8")
        (self.raw / "existing.md").write_text("Existing", encoding="utf-8")

        without = MODULE.build_plan([str(self.raw)], self.vault, recursive=False)
        with_recursive = MODULE.build_plan([str(self.raw)], self.vault, recursive=True)

        self.assertEqual(without[0].action, "invalid")
        self.assertEqual(with_recursive[0].action, "create")
        self.assertEqual(with_recursive[1].action, "skipped")
        self.assertIsNone(with_recursive[1].output)

    def test_empty_recursive_selection_is_invalid(self) -> None:
        items = MODULE.build_plan([str(self.raw)], self.vault, recursive=True)

        self.assertEqual(items[0].action, "invalid")

    def test_unsupported_format_is_reported(self) -> None:
        source = self.raw / "legacy.doc"
        source.write_bytes(b"legacy")

        items = MODULE.build_plan([str(source)], self.vault, recursive=False)

        self.assertEqual(items[0].action, "unsupported")
        self.assertIsNone(items[0].output)

    def test_office_lock_file_is_skipped(self) -> None:
        source = self.raw / "~$draft.docx"
        source.write_bytes(b"lock")

        items = MODULE.build_plan([str(source)], self.vault, recursive=False)

        self.assertEqual(items[0].action, "skipped")

    def test_two_sources_cannot_share_one_sidecar(self) -> None:
        text_source = self.raw / "same.txt"
        html_source = self.raw / "same.html"
        text_source.write_text("Text", encoding="utf-8")
        html_source.write_text("<p>HTML</p>", encoding="utf-8")

        items = MODULE.build_plan(
            [str(text_source), str(html_source)], self.vault, recursive=False
        )

        self.assertEqual([item.action for item in items], ["conflict", "conflict"])
        self.assertEqual(items[0].output, items[1].output)

    def test_docx_conversion_uses_real_engine(self) -> None:
        source = self.raw / "word.docx"
        content_types = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>"""
        relationships = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>"""
        document = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:body><w:p><w:r><w:t>Word conversion marker</w:t></w:r></w:p><w:sectPr/></w:body>
</w:document>"""
        with zipfile.ZipFile(source, "w") as archive:
            archive.writestr("[Content_Types].xml", content_types)
            archive.writestr("_rels/.rels", relationships)
            archive.writestr("word/document.xml", document)

        results, exit_code = MODULE.run_convert(
            [str(source)], self.vault, recursive=False
        )

        self.assertEqual(exit_code, 0, results)
        self.assertIn("Word conversion marker", source.with_suffix(".md").read_text(encoding="utf-8"))

    def test_pptx_conversion_uses_real_engine(self) -> None:
        from pptx import Presentation

        source = self.raw / "slides.pptx"
        deck = Presentation()
        slide = deck.slides.add_slide(deck.slide_layouts[1])
        slide.shapes.title.text = "Slide conversion marker"
        slide.placeholders[1].text = "Slide body"
        deck.save(source)

        results, exit_code = MODULE.run_convert(
            [str(source)], self.vault, recursive=False
        )

        self.assertEqual(exit_code, 0, results)
        output = source.with_suffix(".md").read_text(encoding="utf-8")
        self.assertIn("Slide conversion marker", output)
        self.assertIn("Slide body", output)

    def test_xlsx_conversion_uses_real_engine(self) -> None:
        from openpyxl import Workbook

        source = self.raw / "table.xlsx"
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Data"
        sheet.append(["Metric", "Value"])
        sheet.append(["Conversion marker", 42])
        workbook.save(source)

        results, exit_code = MODULE.run_convert(
            [str(source)], self.vault, recursive=False
        )

        self.assertEqual(exit_code, 0, results)
        output = source.with_suffix(".md").read_text(encoding="utf-8")
        self.assertIn("Conversion marker", output)
        self.assertIn("42", output)

    def test_text_pdf_conversion_uses_real_engine(self) -> None:
        source = self.raw / "document.pdf"
        stream = b"BT /F1 12 Tf 72 720 Td (PDF conversion marker) Tj ET"
        objects = [
            b"<< /Type /Catalog /Pages 2 0 R >>",
            b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>",
            b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
            b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        ]
        pdf = bytearray(b"%PDF-1.4\n")
        offsets = [0]
        for index, obj in enumerate(objects, start=1):
            offsets.append(len(pdf))
            pdf.extend(f"{index} 0 obj\n".encode() + obj + b"\nendobj\n")
        xref = len(pdf)
        pdf.extend(f"xref\n0 {len(objects) + 1}\n".encode())
        pdf.extend(b"0000000000 65535 f \n")
        for offset in offsets[1:]:
            pdf.extend(f"{offset:010d} 00000 n \n".encode())
        pdf.extend(
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode()
        )
        source.write_bytes(pdf)

        results, exit_code = MODULE.run_convert(
            [str(source)], self.vault, recursive=False
        )

        self.assertEqual(exit_code, 0, results)
        self.assertIn("PDF conversion marker", source.with_suffix(".md").read_text(encoding="utf-8"))


if __name__ == "__main__":
    unittest.main()
